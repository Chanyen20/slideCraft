from fastapi import FastAPI, UploadFile, Form
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import FileResponse
import os
import json
import re
from io import BytesIO
from docx import Document
from docx.oxml.ns import qn
from pptx import Presentation
from pptx.util import Inches, Pt
from openai import OpenAI

# 初始化 OpenAI 客戶端
client = OpenAI(api_key="")  # 或用 os.getenv("OPENAI_API_KEY")

# 初始化 FastAPI
app = FastAPI()
app.add_middleware(
    CORSMiddleware,
    allow_origins=["http://localhost:5173"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

UPLOAD_DIR = os.path.join(os.path.dirname(__file__), "uploads")
os.makedirs(UPLOAD_DIR, exist_ok=True)

# GPT 摘要段落為 slide 結構
def summarize_blocks_to_slides(paragraphs, client):
    prompt = (
        "You are given a document with paragraphs. Your task is to group them into slides.\n"
        "Each slide should have a title and 3–5 bullet points derived from the paragraphs.\n"
        "Input paragraphs:\n\n" +
        "\n\n".join(paragraphs) +
        "\n\nNow generate a PowerPoint structure in this format:\n\n"
        "Slide 1:\nTitle: ...\nBullets:\n- ...\n- ...\n\nSlide 2:\nTitle: ...\nBullets:\n- ...\n"
    )

    response = client.chat.completions.create(
        model="gpt-3.5-turbo",
        messages=[{"role": "user", "content": prompt}],
        temperature=0.5,
        max_tokens=1800
    )
    return response.choices[0].message.content

# GPT 回傳的 slide 結構解析
def parse_slide_structure(text):
    slides = []
    current_slide = {"title": "", "bullets": []}
    for line in text.splitlines():
        if line.lower().startswith("slide"):
            if current_slide["title"]:
                slides.append(current_slide)
                current_slide = {"title": "", "bullets": []}
        elif line.lower().startswith("title:"):
            current_slide["title"] = line.split(":", 1)[1].strip()
        elif line.strip().startswith("- "):
            current_slide["bullets"].append(line.strip()[2:])
    if current_slide["title"]:
        slides.append(current_slide)
    return slides

# 將每段原始段落對應到 GPT 產出的 slide
def match_paragraphs_to_slides(paragraphs, slides):
    slide_texts = [" ".join([s["title"]] + s["bullets"]).lower() for s in slides]
    assignments = []
    for para in paragraphs:
        text = para.lower()
        best_match = max(
            range(len(slide_texts)),
            key=lambda i: len(set(re.findall(r'\w+', text)) & set(re.findall(r'\w+', slide_texts[i])))
        )
        assignments.append(best_match)
    return assignments

# 建立 PowerPoint
def generate_fixed_presentation(docx_path, theme, client):
    from pptx.dml.color import RGBColor

    doc = Document(docx_path)
    prs = Presentation()

    blocks = []
    rels = doc.part.rels
    image_rids = {r.rId for r in rels.values() if "image" in r.target_ref}

    for para in doc.paragraphs:
        if para.text.strip():
            blocks.append(('text', para.text.strip()))
        for run in para.runs:
            blips = run._element.xpath(".//a:blip")
            for blip in blips:
                embed = blip.get(qn('r:embed'))
                if embed in image_rids:
                    img_data = rels[embed].target_part.blob
                    blocks.append(('image', img_data))

    seen = set()
    for shape in doc.inline_shapes:
        blip = shape._inline.graphic.graphicData.pic.blipFill.blip
        rId = blip.embed
        if rId in rels and rId not in seen:
            image_data = rels[rId].target_part.blob
            blocks.append(('image', image_data))
            seen.add(rId)

    paragraphs = [b[1] for b in blocks if b[0] == 'text']
    gpt_response = summarize_blocks_to_slides(paragraphs, client)
    slides_data = parse_slide_structure(gpt_response)
    assignments = match_paragraphs_to_slides(paragraphs, slides_data)

    slides = []
    for slide_data in slides_data:
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        title_shape = slide.shapes.title
        title_shape.text = slide_data["title"]
        content_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(4.5))
        tf = content_box.text_frame
        for bullet in slide_data["bullets"]:
            p = tf.add_paragraph()
            p.text = f"• {bullet}"
            p.level = 0
            p.font.size = Pt(28)
        slides.append(slide)

    para_index = 0
    for i, block in enumerate(blocks):
        if block[0] == 'text':
            para_index += 1
        elif block[0] == 'image':
            if para_index > 0:
                slide_idx = assignments[para_index - 1]
                slide = slides[slide_idx]
                image_stream = BytesIO(block[1])
                slide.shapes.add_picture(image_stream, Inches(1), Inches(5.3), height=Inches(2))

    # 可選主題套用（略過背景處理）
    if theme and "text" in theme:
        text_color = RGBColor(
            int(theme["text"][1:3], 16),
            int(theme["text"][3:5], 16),
            int(theme["text"][5:7], 16),
        )
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text_frame") and shape.text_frame:
                    for para in shape.text_frame.paragraphs:
                        para.font.color.rgb = text_color

    pptx_path = os.path.splitext(docx_path)[0] + "_fixed.pptx"
    prs.save(pptx_path)
    return pptx_path

# 上傳處理
@app.post("/upload")
async def upload_file(file: UploadFile, parse_images: bool = Form(False), theme: str = Form("")):
    file_path = os.path.join(UPLOAD_DIR, file.filename)
    with open(file_path, "wb") as f:
        f.write(await file.read())

    theme_data = json.loads(theme) if theme else {}
    pptx_file = generate_fixed_presentation(file_path, theme_data, client)
    return {"pptx_url": f"http://localhost:8000/download/{os.path.basename(pptx_file)}"}

# PPTX 下載
@app.get("/download/{filename}")
async def download_pptx(filename: str):
    pptx_path = os.path.join(UPLOAD_DIR, filename)
    if os.path.exists(pptx_path):
        return FileResponse(pptx_path, filename=filename)
    return {"error": "File not found"}
