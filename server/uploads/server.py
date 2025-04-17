from fastapi import FastAPI, UploadFile, Form
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import FileResponse
import os
import json
from docx import Document
from pptx import Presentation
from pptx.dml.color import RGBColor
import openai
import asyncio
from pptx.util import Inches

# Set OpenAI API key (ensure this key is valid and kept secure)
openai.api_key = ""

# Initialize FastAPI
app = FastAPI()

# Set CORS to allow requests from frontend URL
app.add_middleware(
    CORSMiddleware,
    allow_origins=["http://localhost:5173"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# Set upload directory, create it if it doesn't exist
UPLOAD_DIR = os.path.join(os.path.dirname(__file__), "uploads")
os.makedirs(UPLOAD_DIR, exist_ok=True)

@app.post("/upload")
async def upload_file(file: UploadFile, parse_images: bool = Form(False), theme: str = Form("")):
    """
    Accepts an uploaded Word document and converts it into a PowerPoint presentation.
    """
    # Save the uploaded Word document
    file_path = os.path.join(UPLOAD_DIR, file.filename)
    contents = await file.read()
    with open(file_path, "wb") as f:
        f.write(contents)

    # Parse theme data from frontend (expected in JSON format)
    try:
        theme_data = json.loads(theme)
    except json.JSONDecodeError:
        # Use default theme if parsing fails
        theme_data = {"background": "#FFFFFF", "text": "#000000"}

    # Call async function to generate presentation
    pptx_file = await generate_presentation(file_path, theme_data)

    # Return PPTX download URL
    return {"pptx_url": f"http://localhost:8000/download/{os.path.basename(pptx_file)}"}

async def generate_presentation(docx_path, theme):
    """
    Convert Word document content into a PowerPoint presentation.
    """
    doc = Document(docx_path)
    full_text = "\n\n".join([para.text.strip() for para in doc.paragraphs if para.text.strip()])
    slides_data = await generate_multiple_slides(full_text)

    prs = Presentation()

    # Create slides one by one
    for slide_data in slides_data:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        title_shape = slide.shapes.title
        if title_shape:
            title_shape.text = slide_data.get("title", "")

        # Add a textbox for bullet points
        left = Inches(1)
        top = Inches(2)
        width = Inches(8)
        height = Inches(4)
        content_box = slide.shapes.add_textbox(left, top, width, height)
        text_frame = content_box.text_frame
        text_frame.clear()

        for point in slide_data.get("bullets", []):
            p = text_frame.add_paragraph()
            p.text = point
            p.level = 0

    # Apply the selected theme colors
    apply_theme(prs, theme)

    pptx_filename = os.path.splitext(os.path.basename(docx_path))[0] + ".pptx"
    pptx_path = os.path.join(UPLOAD_DIR, pptx_filename)
    prs.save(pptx_path)

    return pptx_path

async def generate_multiple_slides(full_text):
    """
    Use OpenAI's GPT model to generate slide titles and bullet points from the Word document content.
    """
    prompt = f"""
    You are a helpful assistant that generates a PowerPoint presentation from a Word document.

    Here is the full content:
    """{full_text}"""

    Instructions:
    - Break it into logical sections for a slide deck.
    - For each slide:
    1. Title: 5-7 words starting with an action verb.
    2. Bullet Points: 3-5 clear, concise points.

    Format:
    Slide 1:
    Title: Your Slide Title Here
    Bullets:
    - First bullet point
    - Second bullet point

    Slide 2:
    Title: Your Slide Title Here
    Bullets:
    - First bullet point
    - Second bullet point
    """
    response = await asyncio.to_thread(
        lambda: openai.ChatCompletion.create(
            model="gpt-3.5-turbo",
            messages=[
                {"role": "system", "content": "You generate slide decks from long documents."},
                {"role": "user", "content": prompt}
            ],
            temperature=0.7,
            max_tokens=1500
        )
    )
    output = response.choices[0].message.content
    print("GPT Slide Output:\n", output)
    return parse_multiple_slides(output)

def parse_multiple_slides(output_text):
    """
    Parse GPT response into multiple slide data dictionaries, each containing a title and bullet points.
    """
    slides = []
    current_slide = {}
    for line in output_text.splitlines():
        if line.lower().startswith("slide"):
            if current_slide:
                slides.append(current_slide)
            current_slide = {"title": "", "bullets": []}
        elif line.lower().startswith("title:"):
            current_slide["title"] = line.split(":", 1)[1].strip()
        elif line.strip().startswith("- "):
            current_slide["bullets"].append(line.strip()[2:])
    if current_slide:
        slides.append(current_slide)
    return slides

def apply_theme(prs, theme):
    """
    Apply the selected color theme to all slides.
    """
    try:
        background_color = RGBColor(
            int(theme["background"][1:3], 16),
            int(theme["background"][3:5], 16),
            int(theme["background"][5:7], 16),
        )
        text_color = RGBColor(
            int(theme["text"][1:3], 16),
            int(theme["text"][3:5], 16),
            int(theme["text"][5:7], 16),
        )
    except (KeyError, ValueError):
        background_color = RGBColor(255, 255, 255)
        text_color = RGBColor(0, 0, 0)

    for slide in prs.slides:
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = background_color

        for shape in slide.shapes:
            if hasattr(shape, "text_frame") and shape.text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    if hasattr(paragraph, "font"):
                        paragraph.font.color.rgb = text_color

@app.get("/download/{filename}")
async def download_pptx(filename: str):
    """
    Provide a download link for the generated PowerPoint file.
    """
    pptx_path = os.path.join(UPLOAD_DIR, filename)
    if os.path.exists(pptx_path):
        return FileResponse(pptx_path, filename=filename)
    return {"error": "File not found"}

